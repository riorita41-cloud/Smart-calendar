from pptx import Presentation
from pptx.util import Pt

prs = Presentation()

def add_slide(title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.text = content
    for p in tf.paragraphs:
        p.font.size = Pt(18)

# Слайд 1
slide0 = prs.slides.add_slide(prs.slide_layouts[0])
slide0.shapes.title.text = "Smart Calendar"
slide0.placeholders[1].text = "Интеллектуальный календарь с геймификацией\nдля подготовки к экзаменам\n\n[ВАШЕ ФИО]\n[ГРУППА]"

# Слайды 2-9
add_slide("Проблема и Решение", "❌ ПРОБЛЕМА:\n• Выгорание перед сессией\n• Хаос в планировании\n• Отсутствие мотивации\n\n✅ РЕШЕНИЕ:\n• Авто-генерация расписания\n• Система XP и уровней\n• Единая экосистема для учебы")
add_slide("Конкурентный анализ", "📅 Google Calendar: Только даты, нет мотивации\n Notion/Trello: Сложная настройка, перегруженный UI\n⭐ Smart Calendar: Готов к работе сразу, заточен под учебу, честная система XP")
add_slide("UI/UX Стиль", "• Минимализм: ничего лишнего\n• Профессиональные SVG-иконки (без эмодзи)\n• Чистая типографика\n• Интерфейс не отвлекает от процесса обучения")
add_slide("Киллер-фича №1: Умное расписание", "• Загрузка экзамена + пул вопросов\n• Алгоритм распределяет нагрузку по дням\n• Учитывает дату экзамена и свободное время\n• Результат: готовый план без ручного труда")
add_slide("Киллер-фича №2: Защита XP", "• Начисление за день (+30), задачу (+10), Pomodoro (+5)\n• Поле xpAwarded в БД\n• Проверка статуса ПЕРЕД начислением\n• XP дается строго ОДИН раз за действие")
add_slide("Киллер-фича №3: Pomodoro и задачи", "• Плавающий виджет таймера в фоне\n• Привязка задач к экзаменам\n• Мгновенная обратная связь при выполнении")
add_slide("LIVE DEMO", "Сценарий:\n1. Главная → уровень и активность\n2. Календарь → отметка дня → +30 XP\n3. Повторная отметка → защита от накрутки\n4. Экзамены → сгенерированное расписание\n5. Pomodoro → запуск таймера\n\n⚠️ Демка подготовлена заранее!")
add_slide("Технический стек", "Backend: Symfony (PHP), Doctrine ORM, MySQL\nБезопасность: CSRF-защита, bcrypt хеширование\nDevOps: Миграции БД, Git контроль версий")

# Слайд 10
slide9 = prs.slides.add_slide(prs.slide_layouts[1])
slide9.shapes.title.text = "Спасибо за внимание!"
slide9.placeholders[1].text = "Smart Calendar помогает учиться эффективнее.\n\nПланы: Мобильное приложение, AI-рекомендации\nGitHub: github.com/riorita41-cloud/Smart-calendar\n\n❓ Вопросы?"

prs.save("Smart_Calendar_Presentation.pptx")
print("✅ Презентация создана: Smart_Calendar_Presentation.pptx")